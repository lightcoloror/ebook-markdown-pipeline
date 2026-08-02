from __future__ import annotations

import shutil
import sys
from pathlib import Path
from uuid import uuid4

from PIL import Image
from openpyxl import Workbook
from pptx import Presentation


PROJECT_DIR = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(PROJECT_DIR.parent))
sys.path.insert(0, str(PROJECT_DIR / "scripts"))

from ebook_markdown_pipeline.batch_convert_books import SUPPORTED_FORMATS  # noqa: E402
from ebook_markdown_pipeline.online_document_pipeline import (  # noqa: E402
    IMAGE_EXTENSIONS,
    OnlinePipelineOptions,
    run_online_document_pipeline,
)
from external_wrapper_utils import run_command  # noqa: E402
from generate_quality_fixtures import write_epub_fixture, write_office_fixture  # noqa: E402


def main() -> int:
    root = PROJECT_DIR / f".tmp-online-format-matrix-{uuid4().hex[:8]}"
    try:
        inputs = root / "inputs"
        output = root / "output"
        build_format_fixtures(inputs)
        expected = SUPPORTED_FORMATS | IMAGE_EXTENSIONS
        actual = {path.suffix.lower() for path in inputs.iterdir() if path.is_file()}
        if actual != expected:
            raise AssertionError(f"Fixture coverage mismatch: missing={sorted(expected - actual)}, extra={sorted(actual - expected)}")

        result = run_online_document_pipeline(
            inputs,
            output,
            options=OnlinePipelineOptions(
                provider_mode="fake",
                execute=True,
                request_interval_seconds=0.0,
                max_estimated_cost_usd=0.01,
            ),
        )
        rows = result.get("results") or []
        failures = [row for row in rows if row.get("status") != "ok"]
        outputs = [Path(row["output"]) for row in rows if row.get("output")]
        if result.get("status") != "ok" or failures:
            raise AssertionError(f"Online format matrix failed: {failures}")
        if len(rows) != len(expected) or len(outputs) != len(expected):
            raise AssertionError(f"Expected {len(expected)} completed formats, got rows={len(rows)}, outputs={len(outputs)}")
        if len(set(outputs)) != len(outputs) or not all(path.is_file() for path in outputs):
            raise AssertionError("Online format matrix produced colliding or missing output artifacts.")
        if result.get("safety", {}).get("remote_execution_requested") or any(
            row.get("safety", {}).get("remote_execution_requested") for row in rows
        ):
            raise AssertionError("Fake format matrix must not make remote requests.")
        print(f"Online supported-format matrix passed: {len(expected)} formats.")
        return 0
    finally:
        shutil.rmtree(root, ignore_errors=True)


def build_format_fixtures(root: Path) -> None:
    root.mkdir(parents=True, exist_ok=True)
    markdown = "# Public Format Matrix\n\n## Section\n\nGenerated regression content only.\n"
    (root / "format-md.md").write_text(markdown, encoding="utf-8")
    (root / "format-txt.txt").write_text(markdown, encoding="utf-8")
    (root / "format-html.html").write_text(f"<html><body><h1>Public HTML</h1><p>{markdown}</p></body></html>", encoding="utf-8")
    (root / "format-htm.htm").write_text("<html><body><h1>Public HTM</h1><p>Generated fixture.</p></body></html>", encoding="utf-8")
    (root / "format-csv.csv").write_text("name,value\nalpha,1\nbeta,2\n", encoding="utf-8")
    (root / "format-tsv.tsv").write_text("name\tvalue\nalpha\t1\nbeta\t2\n", encoding="utf-8")
    (root / "format-rtf.rtf").write_text(
        r"{\rtf1\ansi\deff0{\fonttbl{\f0 Arial;}}\viewkind4\uc1\pard\b Public RTF Fixture\b0\par Generated regression content.\par}",
        encoding="ascii",
    )

    write_epub_fixture(root / "format-epub.epub")
    write_office_fixture(root / "format-docx.docx")
    build_pptx(root / "format-pptx.pptx")
    build_xlsx(root / "format-xlsx.xlsx")
    build_pandoc_formats(root, root / "format-md.md")
    build_calibre_formats(root, root / "format-epub.epub")

    shutil.copyfile(PROJECT_DIR / "benchmarks" / "fixtures" / "generated" / "pdf" / "text-layer.pdf", root / "format-pdf.pdf")
    build_image_formats(root)


def build_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Public PPTX Fixture"
    slide.placeholders[1].text = "Generated online-only format regression content."
    presentation.save(path)


def build_xlsx(path: Path) -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Fixture"
    sheet.append(["Metric", "Value"])
    sheet.append(["Alpha", 1])
    sheet.append(["Beta", 2])
    workbook.save(path)


def build_pandoc_formats(root: Path, markdown_source: Path) -> None:
    pandoc = shutil.which("pandoc")
    if not pandoc:
        raise RuntimeError("Pandoc is required for the public online format matrix.")
    for suffix in ("odt", "fb2"):
        output = root / f"format-{suffix}.{suffix}"
        completed = run_command([pandoc, str(markdown_source), "-o", str(output)], root, 60)
        if completed.returncode != 0 or not output.is_file():
            raise RuntimeError(f"Pandoc could not generate {suffix}: {completed.stderr}")


def build_calibre_formats(root: Path, epub_source: Path) -> None:
    ebook_convert = shutil.which("ebook-convert")
    if not ebook_convert:
        raise RuntimeError("Calibre ebook-convert is required for the public online format matrix.")
    for suffix in ("azw3", "mobi"):
        output = root / f"format-{suffix}.{suffix}"
        completed = run_command([ebook_convert, str(epub_source), str(output)], root, 120)
        if completed.returncode != 0 or not output.is_file():
            raise RuntimeError(f"Calibre could not generate {suffix}: {completed.stderr}")
    shutil.copyfile(root / "format-azw3.azw3", root / "format-azw.azw")


def build_image_formats(root: Path) -> None:
    source = PROJECT_DIR / "benchmarks" / "fixtures" / "generated" / "images" / "ocr" / "english.png"
    with Image.open(source) as image:
        rgb = image.convert("RGB")
        for suffix in sorted(IMAGE_EXTENSIONS):
            target = root / f"format-{suffix.lstrip('.')}{suffix}"
            image_format = {".jpg": "JPEG", ".jpeg": "JPEG", ".tif": "TIFF", ".tiff": "TIFF"}.get(suffix)
            if suffix == ".png":
                shutil.copyfile(source, target)
            else:
                rgb.save(target, format=image_format)


if __name__ == "__main__":
    raise SystemExit(main())
